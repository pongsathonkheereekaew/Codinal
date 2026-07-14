"""Verify every shape in a .pptx stays inside the slide canvas.

Usage: python3 check_bounds.py <deck.pptx>
Exits non-zero if any shape overflows (tolerance 0.03in).
"""
import sys
from pptx import Presentation
from pptx.util import Emu

path = sys.argv[1]
prs = Presentation(path)
W = Emu(prs.slide_width).inches
H = Emu(prs.slide_height).inches
tol = 0.03
issues = 0
for i, slide in enumerate(prs.slides, 1):
    for sh in slide.shapes:
        if sh.left is None:
            continue
        l, t = Emu(sh.left).inches, Emu(sh.top).inches
        w = Emu(sh.width).inches if sh.width else 0
        h = Emu(sh.height).inches if sh.height else 0
        probs = []
        if l < -tol or t < -tol:
            probs.append("negative origin")
        if l + w > W + tol:
            probs.append(f"right overflow ({l + w:.2f} > {W:.2f})")
        if t + h > H + tol:
            probs.append(f"bottom overflow ({t + h:.2f} > {H:.2f})")
        if probs:
            issues += 1
            print(f"slide {i}: id={sh.shape_id} {sh.name!r} L={l:.2f} T={t:.2f} "
                  f"W={w:.2f} H={h:.2f} -> {'; '.join(probs)}")
print(f"total issues: {issues}")
sys.exit(1 if issues else 0)
