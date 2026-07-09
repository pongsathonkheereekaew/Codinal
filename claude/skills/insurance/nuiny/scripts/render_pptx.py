import sys, os
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn

SRC = sys.argv[1]
OUT = sys.argv[2]
DPI = 120

FONT = "/System/Library/Fonts/Supplemental/Tahoma.ttf"
FONT_B = "/System/Library/Fonts/Supplemental/Tahoma Bold.ttf"
_font_cache = {}


def font(pt, bold):
    key = (round(pt), bool(bold))
    if key not in _font_cache:
        px = max(6, int(pt * DPI / 72))
        _font_cache[key] = ImageFont.truetype(FONT_B if bold else FONT, px)
    return _font_cache[key]


def emu2px(v):
    return int(Emu(v).inches * DPI)


def solid_fill_color(sh):
    sp = sh._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return None
    sf = spPr.find(qn('a:solidFill'))
    if sf is None:
        return None
    srgb = sf.find(qn('a:srgbClr'))
    if srgb is not None:
        h = srgb.get('val')
        return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))
    return (200, 200, 200)


def cell_fill_color(cell):
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is not None:
        sf = tcPr.find(qn('a:solidFill'))
        if sf is not None:
            srgb = sf.find(qn('a:srgbClr'))
            if srgb is not None:
                h = srgb.get('val')
                return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))
    return (255, 255, 255)


def run_color(r):
    try:
        if r.font.color and r.font.color.type is not None:
            h = str(r.font.color.rgb)
            return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))
    except Exception:
        pass
    return (30, 30, 30)


def wrap_runs(draw, runs, max_w):
    words = []
    for r in runs:
        f = font(r.font.size.pt if r.font.size else 12, r.font.bold)
        c = run_color(r)
        for piece in r.text.split(' '):
            words.append((piece, f, c))
            words.append((' ', f, c))
    if words:
        words.pop()
    lines, cur, cur_w = [], [], 0
    for w, f, c in words:
        ww = draw.textlength(w, font=f)
        if cur and cur_w + ww > max_w and w != ' ':
            lines.append(cur)
            cur, cur_w = [], 0
        if not (not cur and w == ' '):
            cur.append((w, f, c))
            cur_w += ww
    if cur:
        lines.append(cur)
    return lines


ALIGN_C, ALIGN_R = 2, 3


def draw_text_frame(draw, tf, x, y, w, h, vanchor):
    all_lines = []
    for p in tf.paragraphs:
        runs = [r for r in p.runs if r.text]
        if not runs:
            continue
        align = int(p.alignment) if p.alignment is not None else 1
        ls = p.line_spacing if isinstance(p.line_spacing, float) else 1.0
        sb = p.space_before.pt * DPI / 72 if p.space_before else 0
        first = True
        for line in wrap_runs(draw, runs, w):
            lh = max((f.size for _, f, _ in line), default=12) * 1.2 * ls
            all_lines.append((line, lh, align, sb if first else 0))
            first = False
    total_h = sum(lh + sb for _, lh, _, sb in all_lines)
    if vanchor == 3:
        ty = y + max(0, (h - total_h) / 2)
    elif vanchor == 4:
        ty = y + h - total_h
    else:
        ty = y
    for line, lh, align, sb in all_lines:
        ty += sb
        lw = sum(draw.textlength(t, font=f) for t, f, _ in line)
        if align == ALIGN_C:
            tx = x + (w - lw) / 2
        elif align == ALIGN_R:
            tx = x + w - lw
        else:
            tx = x
        for t, f, c in line:
            draw.text((tx, ty), t, font=f, fill=c)
            tx += draw.textlength(t, font=f)
        ty += lh


def render_shape(draw, sh):
    if sh.left is None:
        return
    x, y = emu2px(sh.left), emu2px(sh.top)
    w = emu2px(sh.width or 0)
    h = emu2px(sh.height or 0)
    if sh.shape_type == 6:
        for sub in sh.shapes:
            render_shape(draw, sub)
        return
    if getattr(sh, "has_table", False):
        tbl = sh.table
        col_x = [x]
        for c in tbl.columns:
            col_x.append(col_x[-1] + emu2px(c.width))
        row_y = [y]
        for r in tbl.rows:
            row_y.append(row_y[-1] + emu2px(r.height))
        for ri, cell_row in enumerate(tbl.rows):
            for ci, cell in enumerate(cell_row.cells):
                cx0, cx1 = col_x[ci], col_x[ci + 1]
                cy0, cy1 = row_y[ri], row_y[ri + 1]
                draw.rectangle([cx0, cy0, cx1, cy1], fill=cell_fill_color(cell), outline=(210, 214, 220))
                ml = emu2px(cell.margin_left)
                mr = emu2px(cell.margin_right)
                mt = emu2px(cell.margin_top)
                va = int(cell.vertical_anchor) if cell.vertical_anchor else 1
                draw_text_frame(draw, cell.text_frame, cx0 + ml, cy0 + mt,
                                cx1 - cx0 - ml - mr, cy1 - cy0 - 2 * mt, va)
        return
    fill = solid_fill_color(sh)
    if fill is not None and w > 0 and h > 0:
        draw.rectangle([x, y, x + w, y + h], fill=fill)
    if getattr(sh, "has_text_frame", False):
        tf = sh.text_frame
        va = int(tf.vertical_anchor) if tf.vertical_anchor is not None else 1
        draw_text_frame(draw, tf, x, y, w, h, va)


os.makedirs(OUT, exist_ok=True)
prs = Presentation(SRC)
W = emu2px(prs.slide_width)
H = emu2px(prs.slide_height)
n = 0
for i, slide in enumerate(prs.slides, 1):
    img = Image.new("RGB", (W, H), (10, 25, 47))
    draw = ImageDraw.Draw(img)
    for sh in slide.shapes:
        try:
            render_shape(draw, sh)
        except Exception as e:
            print(f"slide {i} shape {sh.shape_id}: {e}")
    img.save(f"{OUT}/slide{i:02d}.png")
    n = i
print("rendered", n, "slides ->", OUT)
